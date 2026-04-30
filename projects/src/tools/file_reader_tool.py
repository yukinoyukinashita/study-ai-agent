"""
文件读取工具 - 用于读取用户上传的课件文件
"""
import os
from typing import Optional
from langchain.tools import tool


@tool
def read_file(file_path: str, max_length: int = 50000) -> str:
    """
    读取本地文件内容（支持 PDF、Word、PPT、TXT、Markdown 等格式）。
    
    参数:
        file_path: 文件路径（绝对路径或相对路径）
        max_length: 最大读取长度，默认 50000 字符
    
    返回:
        文件的文本内容
    """
    # 处理相对路径
    if not os.path.isabs(file_path):
        workspace_path = os.getenv("COZE_WORKSPACE_PATH", "/workspace/projects")
        file_path = os.path.join(workspace_path, file_path)
    
    if not os.path.exists(file_path):
        return f"文件不存在: {file_path}"
    
    # 根据文件扩展名选择读取方式
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
        elif ext == '.md':
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
        elif ext == '.pdf':
            content = _read_pdf(file_path)
        elif ext in ['.doc', '.docx']:
            content = _read_docx(file_path)
        elif ext in ['.ppt', '.pptx']:
            content = _read_pptx(file_path)
        elif ext in ['.csv']:
            content = _read_csv(file_path)
        elif ext in ['.json']:
            content = _read_json(file_path)
        else:
            # 尝试作为文本文件读取
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            except:
                return f"不支持的文件格式: {ext}"
        
        # 截断过长的内容
        if len(content) > max_length:
            content = content[:max_length] + f"\n\n... (内容已截断，共 {len(content)} 字符)"
        
        return content
        
    except Exception as e:
        return f"读取文件失败: {str(e)}"


def _read_pdf(file_path: str) -> str:
    """读取 PDF 文件"""
    try:
        import pypdf
        reader = pypdf.PdfReader(file_path)
        text_parts = []
        for page in reader.pages:
            text_parts.append(page.extract_text())
        return "\n\n".join(text_parts)
    except ImportError:
        return "PDF 读取需要安装 pypdf 库"


def _read_docx(file_path: str) -> str:
    """读取 Word 文档"""
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n\n".join([p.text for p in doc.paragraphs])
    except ImportError:
        return "Word 文档读取需要安装 python-docx 库"


def _read_pptx(file_path: str) -> str:
    """读取 PowerPoint 文件"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                # 使用 safe way 获取文本
                try:
                    # 尝试获取文本，shape.text 是通用属性
                    text = str(getattr(shape, 'text', ''))
                    if text and text.strip() and text != 'None':
                        text_parts.append(text)
                except Exception:
                    pass
        return "\n\n".join(text_parts)
    except ImportError:
        return "PPT 读取需要安装 python-pptx 库"


def _read_csv(file_path: str) -> str:
    """读取 CSV 文件"""
    try:
        import pandas as pd
        df = pd.read_csv(file_path)
        return df.to_string()
    except ImportError:
        return "CSV 读取需要安装 pandas 库"


def _read_json(file_path: str) -> str:
    """读取 JSON 文件"""
    import json
    with open(file_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    return json.dumps(data, ensure_ascii=False, indent=2)
